import os
import mimetypes
import base64
from io import BytesIO
from openai import OpenAI
from docx import Document
import PyPDF2
from PIL import Image
from pptx import Presentation

# Define folder path
folder_path = "C:/Users/jacky/Desktop/Test files"
files = os.listdir(folder_path)

# Initialize OpenAI client
api_key = os.environ.get("OPENAI_API_KEY")
client = OpenAI()

# Function to resize and encode image in Base64 (without modifying the original file)
def encode_resized_image(file_path, size=(128, 128)):
    try:
        with Image.open(file_path) as img:
            img = img.convert("RGB")  # Ensure RGB format for compatibility

            # Resize only for encoding, preserving the original image
            img_resized = img.resize(size, Image.ANTIALIAS)

            # Save the resized image to a temporary buffer
            buffered = BytesIO()
            img_resized.save(buffered, format="JPEG")  # Use JPEG for efficiency
            base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

            return base64_image
    except Exception as e:
        print(f"Error processing image {file_path}: {e}")
        return base64.b64encode(img.getvalue()).decode("utf-8")

# Function to extract text from files
def extract_text(file_path):
    try:
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            if "text" in mime_type:  # Plain text files
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            elif "msword" in mime_type or file_path.endswith(".docx"):  # Word documents
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                return text if text.strip() else None  # If empty, return None
            elif "pdf" in mime_type:  # PDF files
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
                    return text if text.strip() else None  # If empty, return None
            elif "vnd.ms-powerpoint" in mime_type or file_path.endswith(".pptx"):  # PowerPoint files
                pres = Presentation(file_path)
                text = []
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text) if text else None  # If empty, return None
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return None  # Return None if text extraction fails

# Function to analyze images with OpenAI Vision API
def analyze_image(file_path):
    try:
        # Read and encode image in base64
        with open(file_path, "rb") as image_file:
            base64_image = encode_resized_image(file_path, size=(128, 128))
            #base64_image = base64.b64encode(image_file.read()).decode("utf-8")
        
        # Send the base64-encoded image to OpenAI
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an AI that describes images in a sentence."},
                {
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ],
                },
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error analyzing image {file_path}: {e}")
    return "(No description available)"

# Read files and extract content
file_data = []
existing_folders = [file for file in files if os.path.isdir(os.path.join(folder_path, file))]

for file in files:
    file_path = os.path.join(folder_path, file)
    mime_type, _ = mimetypes.guess_type(file_path)

    if os.path.isdir(file_path):
        continue
    elif mime_type and "image" in mime_type:
        description = analyze_image(file_path)  # Process images
        file_data.append(f"{file}: {description}")
    else:
        text_content = extract_text(file_path)  # Process text-based files
        if text_content:
            file_data.append(f"{file}:\n{text_content[:500]}")  # Limit content to 500 chars
        else:
            file_type = mime_type if mime_type else "unknown"
            description = f"This is a file with name {file}, and file type {file_type}"
            file_data.append(f"{file}: {description}")

print(file_data)



