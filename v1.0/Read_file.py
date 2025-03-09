import os
import mimetypes
import base64
from io import BytesIO
from openai import OpenAI
from docx import Document
import PyPDF2
from PIL import Image
from pptx import Presentation
import shutil

# Initialize OpenAI client
api_key = os.environ.get("OPENAI_API_KEY")
client = OpenAI()

# Function to resize and encode image in Base64 (without modifying the original file)
def encode_resized_image(file_path, size=(128, 128)):
    try:
        with Image.open(file_path) as img:
            img = img.convert("RGB")  # Ensure RGB format for compatibility

            # Resize only for encoding, preserving the original image
            img_resized = img.resize(size)

            # Save the resized image to a temporary buffer
            buffered = BytesIO()
            img_resized.save(buffered, format="JPEG")  # Use JPEG for efficiency
            base64_image = base64.b64encode(buffered.getvalue()).decode("utf-8")

            return base64_image
    except Exception as e:
        print(f"Error processing image {file_path}: {e}")
        return base64.b64encode(img.getvalue()).decode("utf-8")

# Function to extract text from files
def extract_text(file_path):
    try:
        mime_type, _ = mimetypes.guess_type(file_path)
        if mime_type:
            if "text" in mime_type:  # Plain text files
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    return f.read()
            elif "msword" in mime_type or file_path.endswith(".docx"):  # Word documents
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                return text if text.strip() else None  # If empty, return None
            elif "pdf" in mime_type:  # PDF files
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    text = "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
                    return text if text.strip() else None  # If empty, return None
            elif "vnd.ms-powerpoint" in mime_type or file_path.endswith(".pptx"):  # PowerPoint files
                pres = Presentation(file_path)
                text = []
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text.append(shape.text)
                return "\n".join(text) if text else None  # If empty, return None
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
    return None  # Return None if text extraction fails

# Function to analyze images with OpenAI Vision API
def analyze_image(file_path):
    try:
        # Read and encode image in base64
        with open(file_path, "rb") as image_file:
            base64_image = encode_resized_image(file_path, size=(128, 128))
            #base64_image = base64.b64encode(image_file.read()).decode("utf-8")
        
        # Send the base64-encoded image to OpenAI
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": "You are an AI that describes images in a sentence."},
                {
                    "role": "user",
                    "content": [
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{base64_image}"}}
                    ],
                },
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"Error analyzing image {file_path}: {e}")
    return "(No description available)"


def AI_Response(path):
    # Read files and extract content
    files = os.listdir(path)
    file_data = []
    existing_folders = [file for file in files if os.path.isdir(os.path.join(path, file))]

    for file in files:
        file_path = os.path.join(path, file)
        mime_type, _ = mimetypes.guess_type(file_path)

        if os.path.isdir(file_path):
            continue
        elif mime_type and "image" in mime_type:
            description = analyze_image(file_path)  # Process images
            file_data.append(f"{file}: {description}")
        else:
            text_content = extract_text(file_path)  # Process text-based files
            if text_content:
                file_data.append(f"{file}:\n{text_content[:500]}")  # Limit content to 500 chars
            else:
                file_type = mime_type if mime_type else "unknown"
                description = f"This is a file with name {file}, and file type {file_type}"
                file_data.append(f"{file}: {description}")

    # Send request to OpenAI for classification
    completion = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are an assistant designed to optimize and organize directory structures. Based on the provided directory, predict the user's occupation and suggest an improved folder hierarchy that categorizes the files appropriately."},
            {"role": "user", "content": "Generate a list of shell commands to create folders and move the files into appropriate themes based on their content. If the file is hinted to be a folder, then other files can be moved into it if appropriate. There need to be at least two main folders in total (including those that already exist). If there would be a folder containing only one file, instead create a Miscellaneous folder as main folder and group them all in. All files need to be moved into folders. If appropriate, new folders can be created in main folders and move documents in accordingly."},
            {"role": "user", "content": "Make sure that no files are left out."},
            {"role": "user", "content": "Ensure that:  1. Use shutil.move() to move files. 2. Enclose filenames and folder names with spaces in double quotes (`\"...\"`). 3. Use backslashes (`\\`) for file paths instead of forward slashes (`/`). 4. Avoid escaping spaces with backslashes (`\ `). 5.The output should contain only valid Windows commands, separated by `\\n`."},
            {"role": "user", "content": "Here is a list of files with some content: ['cartoon.psd: This is a file with name cartoon.psd, and file type unknown', 'Figure_-3.png: The image shows a plot titled Per Average Log-likelihood Curve, with iterations on the x-axis and average log-likelihood on the y-axis. The curve starts low and increases steadily, indicating improving likelihood with more iterations.', 'Figure_1.png: The image is a scatter plot showing two classes of data points: Class 1 in blue and Class 2 in red, with a title Plot data. The data points appear to be mixed, without clear separability between the two classes.', 'Figure_3.png: The image is a plot titled ''Per Average Log-likelihood Curve,' showing a line graph with iterations on the x-axis and average log-likelihood on the y-axis. The curve, which is red, starts lower and flattens out as it progresses to the right.', 'Figure_7.png: The image shows a plot of scattered data points with two classes: Class 1 in red and Class 2 in blue. The background features contours, likely representing decision boundaries from a classification algorithm.', 'pegin.png: A cartoon penguin wearing an orange jacket and red cap is standing on an icy surface with other penguins nearby.', 'Speech.docx:\nHello everyone. Before we start, I would briefly introduce myself. I am Jacky Cao and I studied in Yanan High school for a year in Grade 10, then I went abroad to Australia to complete my high school education. Here, I will give some information and thoughts on education in Australian and suggestions for whoever who consider going abroad for further university or high school education. \n\nI will first give a big image of foreign education system as well as discussing some differences I spot durin', 'Speech.pptx:\n\nList\nWhy I choose to study abroad\nMy first experience with it\nDifferences and Difficulties\nSuggestions\nWhy\nStressed about Gaokao\nGo overseas to get around it\nIt didnt went well\nI messed up.\nAlien environment\nPoor English\nNo familiar faces\nCry\nThe campus\nMy dorm, shared with 3 other Y10s\nThe name list of my peers in the same house, I have to take it with me all time.\n', 'star.jpg: A silhouette of trees and a building against a colorful sky with hues of pink and purple at dusk or dawn.'], The following folders already exist and should be used instead of creating new ones: [Figure, Cat]"},
            {"role": "assistant", "content": "mkdir Speeches\nmkdir Cartoon\nmkdir Cartoon/Drawing/\nmkdir Miscellaneous\nshutil.move(\"cartoon.psd\", \"Cartoon/\")\nshutil.move(\"Figure_-3.png\", \"Figure/\")\nshutil.move(\"Figure_1.png\", \"Figure/\")\nshutil.move(\"Figure_3.png\", \"Figure/\")\nshutil.move(\"Figure_7.png\", \"Figure/\")\nshutil.move(\"pegin.png\", \"Cartoon/Drawing/\")\nshutil.move(\"Speech.docx\", \"Speeches/\")\nshutil.move(\"Speech.pptx\", \"Speeches/\")\nshutil.move(\"star.jpg\", \"Miscellaneous/\")"},
            {"role": "user", "content": f"Here is a list of files with some content:\n{', '.join(file_data)}, The following folders already exist and should be used instead of creating new ones if applicable: {', '.join(existing_folders)}"},
        ]
    )
    return(completion.choices[0].message.content)


def execute_commands(commands, path):
    # Split the input into individual commands based on 'n'
    print(commands)
    command_list = commands.split('\n')
    os.chdir(path)
    
    for command in command_list:
        # Skip empty strings that might appear due to splitting
        if command.strip():
            if command.strip().startswith('sh'):
                exec(command.strip())
            elif command.strip().startswith ('mkdir'):
                os.system(command)
                
                print(f"Executed: {command}")